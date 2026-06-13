"""產生爛片現形鏡專題簡報。"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from scripts.config import DATA_FILES, OUTPUT_DIR, ROOT

PPT_PATH = OUTPUT_DIR / "專題簡報_爛片現形鏡.pptx"
FONT = "Microsoft JhengHei"


def _font(p, size: int, bold: bool = False) -> None:
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold


def _title_slide(prs, title, sub):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    s.placeholders[1].text = sub
    for sh in (s.shapes.title, s.placeholders[1]):
        for p in sh.text_frame.paragraphs:
            _font(p, 30 if sh == s.shapes.title else 16)


def _bullet(prs, title, lines):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    _font(s.shapes.title.text_frame.paragraphs[0], 26, True)
    body = s.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        _font(p, 18)


def _img(prs, title, path: Path, cap=""):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    tb.text_frame.text = title
    _font(tb.text_frame.paragraphs[0], 24, True)
    if path and str(path) not in (".", "") and path.exists():
        s.shapes.add_picture(str(path), Inches(1), Inches(1), width=Inches(11))
    if cap:
        c = s.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.4))
        c.text_frame.text = cap
        _font(c.text_frame.paragraphs[0], 12)


def build_ppt(members: str = "（請填入組員姓名）") -> Path:
    summary_path = DATA_FILES["analysis_dir"] / "summary_with_ai.json"
    movies = json.loads(summary_path.read_text(encoding="utf-8"))

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _title_slide(prs, "爛片現形鏡", f"電影真實評價分析與雷區預警\n組員：{members}\n網路資料擷取與分析 自選專題")

    _bullet(prs, "研究動機", [
        "電影預告片往往精采，實際觀影卻常踩雷。",
        "PTT 電影版有大量真實鄉民短評，適合做輿情分析。",
        "本專題自動擷取、分析並視覺化，提供避雷參考。",
    ])

    _bullet(prs, "資料擷取", [
        "來源：PTT 電影版（bbs/Movie）",
        "工具：Requests + BeautifulSoup",
        "擷取文章內文與推文（推／噓／→）",
        f"分析電影數：{len(movies)} 部",
    ])

    _bullet(prs, "分析方法", [
        "Jieba 中文斷詞，去除停用詞",
        "統計高頻詞、好評詞、負評詞",
        "依推文標籤（推/噓）計算情緒傾向",
        "WordCloud 產生整體／好評／負評文字雲",
        "Gemini AI 產生 50 字內一句話總結",
    ])

    for m in movies[:3]:
        ai = m.get("ai", {})
        wc = m.get("wordclouds", {})
        wc_path = ROOT / wc.get("all", "") if wc.get("all") else None
        _img(
            prs,
            f"《{m['name']}》— {ai.get('verdict', '分析')}",
            wc_path or Path(""),
            f"{ai.get('summary', '')}｜推/噓：{m.get('push_positive',0)}/{m.get('push_negative',0)}",
        )

    rows = []
    for m in movies:
        ai = m.get("ai", {})
        rows.append(
            f"《{m['name']}》：{ai.get('verdict','—')} — {ai.get('summary','')[:40]}"
        )
    _bullet(prs, "各片 AI 總結", rows)

    _bullet(prs, "Demo 流程", [
        "1. python main.py",
        "2. 開啟 output/index.html",
        "3. 下拉選電影，展示文字雲與 AI 避雷句",
    ])

    _bullet(prs, "限制與倫理", [
        "樣本僅來自 PTT，不代表全體觀眾",
        "研究用途，遵守禮貌爬蟲（延遲請求）",
        "AI 摘要僅供參考，非專業影評",
    ])

    _title_slide(prs, "謝謝聆聽", "Q & A")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPT_PATH))
    print(f"[OK] 簡報：{PPT_PATH}")
    return PPT_PATH


if __name__ == "__main__":
    build_ppt()
